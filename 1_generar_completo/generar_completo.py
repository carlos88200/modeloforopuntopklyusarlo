# import os




# ruta = r'documentos\ENA_2025'
# os.chdir(ruta)
# documentos = os.listdir()
# docs = []
# contar = 0
# for documento in documentos:
#     direccionarchivos = documento  

#     print(f'-----------------------\n ruta {direccionarchivos}')
#     if os.path.isfile(direccionarchivos):
#         pass
        
#     else:
#         print(f'\n', os.listdir(direccionarchivos))
#         direccionarchivossub = os.listdir(direccionarchivos)
#         for arcsub in direccionarchivossub:
#             sub = os.path.join( direccionarchivos,arcsub)
#             print('111111111111111111111')
            
#             if os.path.isfile(sub):
#                 print('archivo')
                
#             else:
#                 print(f'\n', os.listdir(sub))
            
        
# print(f"cuenta = {contar}")
        



# x




















import pdfplumber
import pandas 
from pptx import Presentation
from docx import Document





contar = 0
archivosvalidos = 0

import os
ruta = r'D:\carlos.gonzalez\Documents\documentosForo\documentos\ENA_2025'

#print("Ruta absoluta que se va a recorrer:", ruta,'\n')



carpetas = os.listdir(ruta)
#print(len(carpetas))
carpetasdiccionario = {}

diccionariodelistasdediccionario = {}


for carpeta in carpetas:
    #print(carpeta)
    #carpetasdiccionario[carpeta] = ""
    carpetasdiccionario[carpeta] = []
    rutajoin = ruta+""+"\\"+carpeta
   # print(f'++++{rutajoin}+++\n')
    
    for root, dirs, files in os.walk(rutajoin):
        
        #print(f'Directorio: {root}')
        for nombre_dir in dirs:
            pass
            #print(f'  Subcarpeta: {os.path.join(root, nombre_dir)}\n')
        for nombre_archivo in files:
            #carpetasdiccionario[carpeta] += f"\n\"{nombre_archivo}\"\n"
            contar += 1
            nombre, extencion = os.path.splitext(nombre_archivo)
            #print(f'exttencion: {extencion} y {nombre}')
            #print(f'  Archivo: {os.path.join(root, nombre_archivo)}\n')
            if extencion.lower() == '.pdf':
                archivosvalidos +=1
               
                #carpetasdiccionario[carpeta] += f"\nPdf\n"
                
                with pdfplumber.open(os.path.join(root, nombre_archivo)) as pdf:
                    texto = ""
                    for pd in pdf.pages:
                        texto += pd.extract_text()+'\n'
                        
                    #carpetasdiccionario[carpeta] += texto+'\n'
                    carpetasdiccionario[carpeta].append({'nombre de arcivo':nombre_archivo, 'texto':texto})
                        
                    
            elif extencion.lower() == '.docx':
                archivosvalidos +=1
              #  print('Entro word')
                doc = Document(os.path.join(root, nombre_archivo))
                textol = []
                texto_final = ""
                for d in doc.paragraphs:
                    if d.text.strip():
                        textol.append(d.text.strip())
                        
                        
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                textol.append(cell.text.strip())
                texto_final = "\n".join(textol)
                #print('texto: ',texto)    
                carpetasdiccionario[carpeta].append({'nombre de arcivo':nombre_archivo, 'texto':texto_final})
                    
            elif extencion.lower() == '.xlsx':
                archivosvalidos +=1
                texto = ""
                df = pandas.read_excel(os.path.join(root, nombre_archivo), sheet_name=None)
                #carpetasdiccionario[carpeta] += f"\nxlsx\n"

                json = ""
                for nombre, items in df.items():
                    #carpetasdiccionario[carpeta] +='Hoja'  
                    columnaspd = items.columns
                    for index, row in items.iterrows():
                        json += "{\n"
                        for colpd in columnaspd:
                            json += f"{colpd}: {row[colpd]},\n"
                        json += "\n}\n"
                        
                    #texto = items.to_string(index=False)
                    #print('-------------------Joha', texto)
                #carpetasdiccionario[carpeta] += json+'\n'
                carpetasdiccionario[carpeta].append({'nombre de arcivo':nombre_archivo, 'texto':texto})


            elif extencion.lower() == '.pptx':
                archivosvalidos +=1
                texto= ""
                pre = Presentation(os.path.join(root, nombre_archivo))
                texto += 'presentacion'
                for num_slide, slide in enumerate(pre.slides, start=1):
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texto += shape.text+'\n'
                #carpetasdiccionario[carpeta] += texto+'\n'   
                carpetasdiccionario[carpeta].append({'nombre de arcivo':nombre_archivo, 'texto':texto})
                    
            
            
        
        
print('\n\tFinal: ', contar)
print('\n\tFinal: ', archivosvalidos)





#print(carpetasdiccionario['1_AGENDA_actualizado_liberado'])

dataframes = {}
for clave, valor in carpetasdiccionario.items():
    
    #print('----',clave,'\n')
    
    
    dataframes[clave] = pandas.DataFrame(valor)
    # for val in valor:
    #     print(val)
        
print('se acabo \n\n\n')
    
    
    
    
    
    
    
    
    
for clave, volar in dataframes.items():
    #print(type(volar))
    volar.to_excel(f'{clave}.xlsx', index=False)
    

completo = pandas.concat(dataframes.values(), ignore_index=True)
completo.to_excel('Completo.xlsx', index=False)
completo.to_csv('Completo.txt', sep=' ', index=False)

# r = r"D:\carlos.gonzalez\Documents\documentosForo\documentos\ENA_2025\1_AGENDA_actualizado_liberado\AGENDA_Cap_ENT y JC_ENA2025_25092025.docx"
# doc = Document(r)
# for i in doc.paragraphs:
#     print('-'+i.text)
    
    
    
    
# ruta = r"D:\carlos.gonzalez\Documents\documentosForo\documentos\ENA_2025\1_AGENDA_actualizado_liberado\AGENDA_Cap_ENT y JC_ENA2025_25092025.docx"
# doc = Document(ruta)

# # Párrafos "libres"
# for p in doc.paragraphs:
#     if p.text.strip():
#         print('[Párrafo] ' + p.text)

# # Párrafos dentro de tablas
# for table in doc.tables:
#     for row in table.rows:
#         for cell in row.cells:
#             texto_celda = cell.text.strip()
#             if texto_celda:
#                 print('[Tabla] ' + texto_celda)


